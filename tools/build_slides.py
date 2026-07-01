#!/usr/bin/env python3
"""Build a concise presentation deck from REPORT.md.

Generates REPORT_slides.pptx — ~23 slides, zh-TW + English terms, results-first.
Sections: Introduction / Methodology / Results / Findings / Discussion.
Embeds key figures from figures/out/. Re-run after editing this script.

    /tmp/pptenv/bin/python tools/build_slides.py      # (needs python-pptx)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIG = os.path.join(HERE, "figures", "out")
OUT = os.path.join(HERE, "REPORT_slides.pptx")

# ---- palette --------------------------------------------------------------
PRIMARY = RGBColor(0x16, 0x3A, 0x5F)   # deep blue (titles)
ACCENT  = RGBColor(0x2E, 0x86, 0xC1)   # blue (rules, tags)
GOOD    = RGBColor(0x1E, 0x8E, 0x3E)   # green (wins)
BAD     = RGBColor(0xC0, 0x39, 0x2B)   # red (loses)
GRAY    = RGBColor(0x60, 0x60, 0x60)
LIGHT   = RGBColor(0xED, 0xF2, 0xF7)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

_total = 0  # slide counter (filled at the end)


def add_slide(section=None):
    s = prs.slides.add_slide(BLANK)
    return s


def set_para(p, text, size, color, bold=False, align=PP_ALIGN.LEFT, level=0):
    p.level = level
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return p


def add_header(slide, title, section=None):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.9))
    tf = box.text_frame
    tf.word_wrap = True
    set_para(tf.paragraphs[0], title, 27, PRIMARY, bold=True)
    # accent rule
    line = slide.shapes.add_shape(1, Inches(0.62), Inches(1.18), Inches(3.2), Pt(3))
    line.fill.solid(); line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    line.shadow.inherit = False
    if section:
        tag = slide.shapes.add_textbox(Inches(9.9), Inches(0.42), Inches(2.8), Inches(0.4))
        set_para(tag.text_frame.paragraphs[0], section, 12, ACCENT, bold=True,
                 align=PP_ALIGN.RIGHT)


def add_footer(slide, idx):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(7.05), Inches(12.1), Inches(0.35))
    set_para(box.text_frame.paragraphs[0],
             f"SQLite Cold-Start Prefetch · Preprocessing Cost-Accounting        {idx}",
             9, GRAY)


def body_frame(slide, top=1.45, height=5.3, left=0.7, width=11.9):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    return tf


def bullets(slide, items):
    """items: list of (text, level, color, bold). First para reuses paragraph[0]."""
    tf = body_frame(slide)
    first = True
    for text, level, color, bold in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        size = 18 if level == 0 else 15
        prefix = "● " if level == 0 else "–  "
        # bullet glyph via run
        rb = p.add_run(); rb.text = prefix
        rb.font.size = Pt(size); rb.font.color.rgb = (ACCENT if level == 0 else GRAY)
        rb.font.bold = False
        rt = p.add_run(); rt.text = text
        rt.font.size = Pt(size); rt.font.color.rgb = color; rt.font.bold = bold
        p.space_after = Pt(7 if level == 0 else 4)


def B(text, level=0, color=None, bold=False):
    return (text, level, color or RGBColor(0x22, 0x22, 0x22), bold)


def add_figure(slide, fname, caption, max_w=10.6, max_h=4.65, top=1.5):
    path = os.path.join(FIG, fname)
    pic = slide.shapes.add_picture(path, 0, 0)
    nw, nh = pic.width, pic.height
    scale = min(Inches(max_w) / nw, Inches(max_h) / nh)
    pic.width = int(nw * scale); pic.height = int(nh * scale)
    pic.left = int((SW - pic.width) / 2)
    pic.top = Inches(top)
    cap = slide.shapes.add_textbox(Inches(0.7), Inches(top) + pic.height + Inches(0.08),
                                   Inches(11.9), Inches(0.7))
    cap.text_frame.word_wrap = True
    set_para(cap.text_frame.paragraphs[0], caption, 12.5, GRAY, align=PP_ALIGN.CENTER)


def add_table(slide, headers, rows, top=1.6, col_w=None, fontsz=13,
              cell_colors=None):
    nrows, ncols = len(rows) + 1, len(headers)
    total_w = 12.0
    width = Inches(total_w)
    height = Inches(0.45 * nrows)
    gtab = slide.shapes.add_table(nrows, ncols, Inches(0.66), Inches(top),
                                  width, height).table
    if col_w:
        for j, w in enumerate(col_w):
            gtab.columns[j].width = Inches(w)
    # header
    for j, h in enumerate(headers):
        c = gtab.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = PRIMARY
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = c.text_frame.paragraphs[0]
        set_para(p, h, fontsz, WHITE, bold=True, align=PP_ALIGN.CENTER)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = gtab.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = LIGHT if i % 2 else WHITE
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            color = RGBColor(0x22, 0x22, 0x22)
            if cell_colors and (i, j) in cell_colors:
                color = cell_colors[(i, j)]
            p = c.text_frame.paragraphs[0]
            set_para(p, val, fontsz, color, bold=(j == 0),
                     align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)


# ===========================================================================
# 1 · Title
# ===========================================================================
s = add_slide()
band = s.shapes.add_shape(1, 0, Inches(2.35), SW, Inches(2.7))
band.fill.solid(); band.fill.fore_color.rgb = PRIMARY
band.line.fill.background(); band.shadow.inherit = False
t = s.shapes.add_textbox(Inches(0.8), Inches(2.55), Inches(11.7), Inches(1.5))
t.text_frame.word_wrap = True
set_para(t.text_frame.paragraphs[0],
         "SQLite Cold-Start Prefetch 的 Preprocessing Cost-Accounting",
         34, WHITE, bold=True)
sub = s.shapes.add_textbox(Inches(0.8), Inches(4.05), Inches(11.7), Inches(0.8))
set_para(sub.text_frame.paragraphs[0],
         "為何 first-query 改善 ≠ end-to-end 加速", 20, RGBColor(0xCF, 0xE3, 0xF3))
foot = s.shapes.add_textbox(Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.8))
set_para(foot.text_frame.paragraphs[0],
         "範圍：commodity x86 桌機 + NVMe｜實驗平台 Ryzen 9950X｜全 cell cold_pct=0",
         13, GRAY)

# ===========================================================================
# 2 · Agenda
# ===========================================================================
s = add_slide(); add_header(s, "大綱 Agenda")
bullets(s, [
    B("Introduction — 問題、兩大挑戰、部署模型", 0, PRIMARY, True),
    B("Methodology — 量測協定、兩個部署模型、selection vs delivery", 0, PRIMARY, True),
    B("Results — 核心 trade-off、ablation、競爭性 baseline", 0, PRIMARY, True),
    B("Findings — 五維 robustness、關鍵發現", 0, PRIMARY, True),
    B("Discussion — 實務建議、限制、結論", 0, PRIMARY, True),
])

# ===========================================================================
# INTRODUCTION
# ===========================================================================
SEC = "Introduction"
s = add_slide(); add_header(s, "問題：SQLite cold-start 第一筆查詢很慢", SEC)
bullets(s, [
    B("SQLite 是部署最廣的 DB（估 >1 兆個 in use）；每次 app 喚醒、裝置恢復，使用者感知的「第一筆查詢延遲」由 cold-start 決定"),
    B("cold-start = OS page cache 為空；B+tree 每筆 query 要走 root→leaf，每個 page miss = 5–100 µs random I/O"),
    B("cold first-query 比 warm 慢 200 倍以上；baseline 落在 ~529–1096 µs（A 529 / B 760 / C 1096 µs）", 1),
    B("關鍵結構事實：interior page 只占 0.35%（92 個 / 368 KB）卻是每筆 query 必經 → targeted prefetch 的著力點", 0, GOOD, True),
])
s = add_slide(); add_header(s, "兩個尚未被同時解決的挑戰", SEC)
bullets(s, [
    B("Targeting：OS 與應用層都看不到 SQLite B+tree 的 page-type 結構 → 盲目 prefetch 浪費 I/O 與 page reclaim", 0, PRIMARY, True),
    B("Cost-accounting：既有 prefetch 多只優化 first-query，未把 prefetch 自身的 preprocessing 計入 end-to-end", 0, PRIMARY, True),
    B("造成「first-query 改善」被誤讀為「真實 cold-start 變快」的系統性誤導", 1, BAD, False),
    B("本研究同時處理這兩點：page-type / access-frequency targeting ＋ 把 preprocessing 拆到 OS-syscall 粒度計入 critical path"),
])
s = add_slide(); add_header(s, "部署形態：warm process, cold data", SEC)
bullets(s, [
    B("現代 cold-start 多是「process 還在、資料已冷」：serverless keep-alive、microservice 重啟、app warm start（有文獻佐證）"),
    B("兩個部署模型：", 0, PRIMARY, True),
    B("warm-process / integrated：app 已在跑、重用既有 handle、不付冷 open（本研究主張）", 1, GOOD, True),
    B("standalone warmer：另起 process、需付冷 open（較悲觀的對照）", 1, GRAY, False),
    B("範圍界定：所有量測在 commodity x86 桌機 + NVMe；mobile / IoT 是 motivation、未在其上量測", 0, BAD, False),
])
s = add_slide(); add_header(s, "研究問題與貢獻", SEC)
bullets(s, [
    B("RQ1 targeting／RQ2 cost-accounting／RQ3 selection vs delivery／RQ4 robustness", 0, PRIMARY, True),
    B("C1 type-aware layout rewriter — 探索性負面結果（預設不部署）", 1),
    B("C2 access-pattern frugality — 極少 syscall 取得 first-query −26~81%", 1),
    B("C3 preprocessing cost-accounting 框架（兩部署模型）— 核心方法貢獻", 1, GOOD, True),
    B("C4 五維 robustness 驗證（churn / RAM / cadence / 10-seed / 1 GiB）", 1),
])

# ===========================================================================
# METHODOLOGY
# ===========================================================================
SEC = "Methodology"
s = add_slide(); add_header(s, "測試 DB 與 workloads", SEC)
add_table(s,
    ["Workload", "存取型態", "典型場景"],
    [["A", "Zipfian point read（熱門集中）", "App 首頁、常開聯絡人"],
     ["B", "Uniform random（隨機讀）", "隨機抽樣、爬蟲"],
     ["C", "High-key（查 file tail 最新）", "剛收到的訊息 / 剛拍的照片"],
     ["D", "Write-heavy churn", "製造 layout 漂移（§6.2.1）"]],
    top=1.55, col_w=[1.6, 5.6, 4.8])
tf = body_frame(s, top=4.5, height=2.0)
set_para(tf.paragraphs[0],
         "Reference DB：600,000 row · 102 MB · 26,331 page · interior 92 個（0.35% / 368 KB）。"
         "三種 layout（1a 原始 / 1b VACUUM / 1c type-aware）共用同一份內容、僅 page 物理排列不同。",
         14, GRAY)
s = add_slide(); add_header(s, "Cold-start 量測協定", SEC)
bullets(s, [
    B("「warm process, cold data」：process 層（handle / mmap / CPU cache）保持暖，資料層每次量測前歸零"),
    B("每個 cell 的把關：", 0, PRIMARY, True),
    B("全機 drop-caches → harness `--verify-hotset`（mincore，cold_pct≈0；>1% 剔除）", 1),
    B("`majflt > 0` 驗證每筆 I/O 真的打到 disk（非模擬器）", 1),
    B("計時前先把 CPU busy-spin 拉到滿頻，消除「最快 cell 受 freq ramp 懲罰最重」的偏差", 1),
    B("warm vs full-cold 偏差估 1–3 µs（< baseline 的 1%）"),
])
s = add_slide(); add_header(s, "兩個部署模型的 end-to-end", SEC)
bullets(s, [
    B("preprocessing 拆成兩個 term：", 0, PRIMARY, True),
    B("open(db)：冷開 DB ~200 µs，per-layout 常數、與策略無關（common-mode）", 1),
    B("deliver：逐頁 madvise / pread，隨 hotset 大小（小策略 ~70–200 µs、2f 可達 ~7 ms）", 1),
    B("e2e_std = open + deliver + first-query（standalone）", 0, GRAY, True),
    B("e2e_warm = deliver + first-query（warm-process，本研究主張）", 0, GOOD, True),
    B("open 只決定兩模型之差，不決定 warm 模型內 prefetch 對 baseline 的勝負"),
])
s = add_slide(); add_header(s, "selection vs delivery（兩種模式）", SEC)
bullets(s, [
    B("`madvise(WILLNEED)` 是 async hint、不保證載到 → first-query 同時受「選對頁」與「載得及」影響"),
    B("用兩種 delivery 模式把這兩件事拆開：", 0, PRIMARY, True),
    B("pread 模式（oracle，同步保證載入）= 「hotset 選對了嗎」的 first-query 下界", 1, GOOD, True),
    B("async 模式（實務）= 「實際拿到多少」（配 delivery_pct）", 1),
    B("delivery loss = fq_async − fq_pread；實測：per-page warmer 在 sleep=0 已 100% 落地", 0, GRAY, False),
])
s = add_slide(); add_header(s, "Strategies 與統計方法", SEC)
bullets(s, [
    B("Layout：1a 原始 / 1b VACUUM / 1c type-aware（interior 集中至 file head）"),
    B("Prefetch：2a–2c 看結構（layers_N）／ 2d–2e access-pattern（interior + top-K 熱 leaf）／ 2f 抄整份 cache"),
    B("統計：每 cell n=10、報 median（對長尾 random-I/O 比 mean 穩健）", 0, PRIMARY, True),
    B("10-seed bootstrap 95% CI 量「換一條同分佈不同抽樣的 workload」的不確定性", 1),
    B("判定：CI 不跨 0 → robust；跨 0 但 ≥7/10 同號 → directional；否則 tie", 1),
])

# ===========================================================================
# RESULTS
# ===========================================================================
SEC = "Results"
s = add_slide(); add_header(s, "每 workload：first-query 最低 ≠ e2e 最佳", SEC)
add_table(s,
    ["Workload", "first-q 最低", "first-q", "e2e best（warm-process）"],
    [["A Zipfian", "2f_slru", "127 µs (−76%)", "layers_5 / 2d / 2e_K10  −7~9%"],
     ["B uniform", "2f_slru", "128 µs (−83%)", "2d / 2e_K10  −29~34%"],
     ["C file-tail", "2f_slru", "123 µs (−89%)", "2e_K10  −73%（291 µs，全矩陣最佳）"]],
    top=1.7, col_w=[2.4, 2.0, 2.4, 5.2],
    cell_colors={(3, 3): GOOD})
tf = body_frame(s, top=4.6, height=1.8)
set_para(tf.paragraphs[0],
         "2f_slru（載整份 working set）三 workload first-query 都最低，但其 deliver 0.8–7 ms 使 e2e 多半輸；"
         "warm-process 下小而準的 targeted prefetch 三 workload e2e 都改善。",
         14, GRAY)
s = add_slide(); add_header(s, "核心觀察：first-query ≠ end-to-end", SEC)
add_figure(s, "14_strategy_endtoend_stacked.png",
           "圖 14：end-to-end cold start（stacked：first-q + deliver + 灰=冷 open）。2f 兩模型都遠超 baseline（A +1248%）；"
           "warm-process 下 targeted prefetch 三 workload 都在 baseline 之下。")
s = add_slide(); add_header(s, "Workload C：小而準最有效益", SEC)
add_table(s,
    ["做法", "first-q", "warm e2e", "跨 10 seed (CI)"],
    [["載全部 interior (layers_92)", "−37%", "−20%", "−21% robust"],
     ["只載用過的 interior (2d)", "−38%", "−31%", "−36% robust"],
     ["+ 最熱 10 個 leaf (2e_K10)", "−81%", "−73%（291 µs）", "−70% [−72,−69] robust"]],
    top=1.7, col_w=[5.0, 2.0, 2.6, 2.4],
    cell_colors={(3, 1): GOOD, (3, 2): GOOD, (3, 3): GOOD})
tf = body_frame(s, top=4.5, height=1.8)
set_para(tf.paragraphs[0],
         "三者 deliver 都小（~70–200 µs）→ e2e 由 first-query 決定；加載少量熱 leaf（2e_K10）"
         "把 first-q 壓到 −81%、warm e2e −73%，是全研究最穩健的勝負。",
         14, GRAY)
s = add_slide(); add_header(s, "Ablation：贏在 access-frequency，不是 page-type", SEC)
add_figure(s, "17_lever_ablation.png",
           "圖 17：三槓桿 ablation。同型別、同張數下，隨機 10 葉只 −2%（灰），照頻率挑的 10 葉 −40%（綠）"
           "→ C 的勝利是 access-frequency 訊號；page-type（interior）撐起 uniform B。")
s = add_slide(); add_header(s, "競爭性 baseline：不是贏稻草人", SEC)
add_figure(s, "18_competitive_baseline.png",
           "圖 18：對打調校過的 frequency-ranked partial dump（2f_topN，純頻率、零 page-type）。"
           "broad A/B 追平（page-type 非必要）；narrow C 仍勝；full dump 因 deliver 過重在 A/B 爆 +730~762%。")

# ===========================================================================
# FINDINGS
# ===========================================================================
SEC = "Findings"
s = add_slide(); add_header(s, "Robustness 1：10-seed 交叉驗證", SEC)
bullets(s, [
    B("access-pattern targeted prefetch 三 workload 的 warm e2e 皆 robust（CI 不跨 0、≥9/10 同號）：", 0, PRIMARY, True),
    B("C 2e_K10 −70% [−72,−69]、A 2e_K10 −36% [−50,−23]、B 2d −25% [−32,−16]", 1, GOOD, True),
    B("structural layers_5 在 A/B 落在雜訊內（tie / directional、CI 跨 0）→ 不可恃", 0, BAD, False),
    B("方法學教訓：單一 workload 點估計會誤導 — 原始 A 抽到「便宜的第一筆查詢」，"
      "同時低估 targeted、高估 layers_5；只有多 seed 才看得出來"),
])
s = add_slide(); add_header(s, "Robustness 2：RAM 壓力 / churn / DB 放大", SEC)
add_figure(s, "16_ram_pressure_sweep.png",
           "圖 16：cap 壓到 working set 以下。targeted（≤2 MB hotset）全程 delivery_pct=100% → first-q 平；"
           "2f_slru（17.7 MB dump）delivery 隨 cap 線性塌、first-q 跳回 baseline（all-or-nothing）。",
           max_h=4.2)
tf = body_frame(s, top=6.0, height=0.9)
set_para(tf.paragraphs[0],
         "另：50k write churn 下 static hotset 不 decay；DB 放大 10× 到 1 GiB，targeted first-q 效益 size-robust，"
         "cache-dump 的 deliver 陷阱隨 DB 惡化。", 12.5, GRAY)
s = add_slide(); add_header(s, "關鍵發現總結", SEC)
bullets(s, [
    B("first-query ≠ end-to-end：2f first-q 最低（−76~89%）但 e2e 多半輸（deliver 太重）", 0, PRIMARY, True),
    B("部署模型決定勝負：warm-process 下 targeted prefetch 三 workload e2e 都贏", 0, PRIMARY, True),
    B("效益主要來自 access-frequency + 小 footprint；page-type 在 narrow workload 加值（保 path coverage）"),
    B("旗艦結果：C × 2e_K10 warm e2e −73%（單一 workload）／ −70% [−72,−69]（跨 10 seed）", 0, GOOD, True),
])

# ===========================================================================
# DISCUSSION
# ===========================================================================
SEC = "Discussion"
s = add_slide(); add_header(s, "實務建議", SEC)
add_table(s,
    ["場景", "建議做法", "first-q", "warm e2e"],
    [["慢 workload（查 tail / churn）", "2e_K10：interior + ~10 熱 leaf", "−81%", "−73%（最佳）"],
     ["uniform 隨機讀", "structural layers_5 / 2d", "−42~43%", "−29~34%"],
     ["快 workload（熱門集中）", "access-pattern 2d / 2e（避免單用 layers_5）", "−22~26%", "−25~36% robust"],
     ["Batch / 平均 latency", "2f SLRU（抄整份 cache）", "−76~89%", "多半不具優勢"],
     ["多 process 共享 DB", "shared cache + background warmer", "—", "cadence ≤ query 間隔"]],
    top=1.65, col_w=[3.2, 4.6, 1.9, 2.3], fontsz=12.5,
    cell_colors={(1, 3): GOOD, (4, 3): BAD})
s = add_slide(); add_header(s, "限制與範圍", SEC)
bullets(s, [
    B("單機 / 單 kernel / 單 SSD / 單 DB：絕對 µs 跨 session 漂 30–70% → 用 2f_slru 當錨點、跨批只比相對量"),
    B("ra=128 單值；async/pread gap 的「readahead 養熱」成因標為 labeled conjecture（無 root 沒掃 ra）", 1),
    B("mobile / ARM / UFS 未量測（motivation、非 evaluated platform）；C 的 sub-WS RAM-robustness 為演繹", 1),
    B("Future work：NVMe stream directives 物理隔離（FEMU）、ARM/UFS SBC 重跑關鍵 cell、strict cold-start", 0, ACCENT, False),
])
s = add_slide(); add_header(s, "結論 Conclusion", SEC)
bullets(s, [
    B("first-query 改善 ≠ end-to-end 加速 — 把 prefetch preprocessing 計入 critical path 才看得到真相", 0, PRIMARY, True),
    B("warm-process 下 access-pattern targeted prefetch 跨 seed robust（C −70%、A −36%、B −25%）", 0, GOOD, True),
    B("cache-dump（2f）first-query 看似最低，但 deliver 0.8–7 ms 使 e2e 慢一個量級 — 此 trade-off 長期被忽略"),
    B("核心貢獻：SQLite cold-start 上 OS-syscall 粒度的 preprocessing cost-accounting ＋ 兩部署模型對齊"),
])

# ---- finalize: footers + page numbers ------------------------------------
_total = len(prs.slides._sldIdLst)
for i, slide in enumerate(prs.slides, start=1):
    if i == 1:
        continue
    add_footer(slide, f"{i} / {_total}")

prs.save(OUT)
print(f"wrote {OUT}  ({_total} slides)")
